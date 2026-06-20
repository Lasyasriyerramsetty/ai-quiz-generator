import os
import json
import re
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 25 * 1024 * 1024  # 25 MB
app.config['UPLOAD_FOLDER'] = 'uploads'

os.makedirs('uploads', exist_ok=True)

API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
IS_OPENROUTER = API_KEY.startswith("sk-or-")

if IS_OPENROUTER:
    from openai import OpenAI
    llm = OpenAI(base_url="https://openrouter.ai/api/v1", api_key=API_KEY)
    MODEL = "openai/gpt-oss-120b:free"
else:
    import anthropic
    llm = anthropic.Anthropic(api_key=API_KEY)
    MODEL = "claude-haiku-4-5-20251001"


def call_llm(prompt: str, max_tokens: int = 4096) -> str:
    if IS_OPENROUTER:
        resp = llm.chat.completions.create(
            model=MODEL,
            messages=[{"role": "user", "content": prompt}],
            max_tokens=max_tokens,
        )
        return resp.choices[0].message.content
    else:
        resp = llm.messages.create(
            model=MODEL,
            max_tokens=max_tokens,
            messages=[{"role": "user", "content": prompt}],
        )
        return resp.content[0].text


DIFFICULTY_DESC = {
    'Simple': 'basic recall and definition questions that test memory of key terms and facts',
    'Medium': 'balanced recall and scenario questions requiring understanding and application of concepts',
    'Complex': 'advanced analytical and scenario-based questions requiring deep understanding and critical thinking',
}


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    if not file.filename or not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Only .pptx files are accepted'}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    try:
        prs = Presentation(filepath)
        slides_text = []
        word_count = 0

        for i, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_content.append(shape.text.strip())
            if slide_content:
                text = ' '.join(slide_content)
                slides_text.append({'slide': i + 1, 'text': text})
                word_count += len(text.split())

        total_slides = len(prs.slides)
    finally:
        os.remove(filepath)

    if not slides_text:
        return jsonify({'error': 'No readable text found in the presentation'}), 400

    return jsonify({
        'filename': filename,
        'slide_count': total_slides,
        'word_count': word_count,
        'slides': slides_text,
    })


@app.route('/api/generate', methods=['POST'])
def generate_quiz():
    data = request.json or {}
    slides = data.get('slides', [])
    num_questions = max(5, min(30, int(data.get('num_questions', 10))))
    difficulty = data.get('difficulty', 'Medium')

    if not slides:
        return jsonify({'error': 'No slide content provided'}), 400

    full_text = '\n\n'.join([f"Slide {s['slide']}: {s['text']}" for s in slides])
    desc = DIFFICULTY_DESC.get(difficulty, DIFFICULTY_DESC['Medium'])

    prompt = f"""You are an expert quiz creator. Based on the PowerPoint slide content below, generate exactly {num_questions} multiple-choice questions at {difficulty} difficulty ({desc}).

SLIDE CONTENT:
{full_text}

STRICT RULES:
1. Exactly {num_questions} questions — no more, no less.
2. Each question has exactly 4 options: A, B, C, D.
3. Exactly one option is correct; the other three are plausible distractors.
4. Distribute questions evenly across different topics/slides.
5. Match {difficulty} difficulty precisely.
6. Make all options similar in length so the correct answer isn't obvious.

Return ONLY a valid JSON array — no markdown fences, no extra text:
[
  {{
    "id": 1,
    "question": "Question text?",
    "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
    "correct": "B",
    "topic": "Short topic label"
  }}
]"""

    try:
        raw = call_llm(prompt, max_tokens=4096)
        raw = raw.strip()
        raw = re.sub(r'^```[a-z]*\s*', '', raw)
        raw = re.sub(r'\s*```$', '', raw)
        questions = json.loads(raw)
        return jsonify({'questions': questions})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/feedback', methods=['POST'])
def get_feedback():
    data = request.json or {}
    wrong_answers = data.get('wrong_answers', [])

    if not wrong_answers:
        return jsonify({'feedback': []})

    feedback_list = []
    for item in wrong_answers:
        prompt = f"""A student answered a quiz question incorrectly. Write a clear, concise explanation (2-3 sentences) of why their answer is wrong and why the correct answer is right.

Question: {item['question']}
Student answered: {item['user_answer']} — "{item['user_option_text']}"
Correct answer: {item['correct_answer']} — "{item['correct_option_text']}"

Reply with only the explanation. No labels, no prefixes."""

        try:
            explanation = call_llm(prompt, max_tokens=256)
            feedback_list.append({
                'question_id': item['question_id'],
                'explanation': explanation.strip(),
            })
        except Exception as e:
            feedback_list.append({
                'question_id': item['question_id'],
                'explanation': f'Could not generate explanation: {e}',
            })

    return jsonify({'feedback': feedback_list})


if __name__ == '__main__':
    app.run(debug=True, port=5000)
